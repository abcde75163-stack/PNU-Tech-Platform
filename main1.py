# =========================================================
# Colab 통합(A안) – PDF + PPTX + DOCX 업로드
# ✅ 최종 통합본 vFinal-Rewrite-Optimized-Hybrid-v9 (Streamlit 버전)
# [필요 라이브러리 설치]
# pip install streamlit pypdf openai python-pptx pillow pymupdf requests numpy python-docx
# =========================================================

import os, re, io, json, tempfile, traceback, base64, hashlib, time, inspect
from datetime import datetime
from concurrent.futures import ThreadPoolExecutor, as_completed

import requests
import numpy as np
import streamlit as st
from pypdf import PdfReader
import docx
from pptx import Presentation
from pptx.util import Cm
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from PIL import Image, ImageChops
import fitz  # PyMuPDF

import openai
from openai import OpenAI

# =========================================================
# 0) API KEY (환경변수만 사용)
# =========================================================
os.environ["OPENAI_API_KEY"] = st.secrets["OPENAI_API_KEY"]
OPENAI_API_KEY = os.environ.get("OPENAI_API_KEY", "").strip()
if not OPENAI_API_KEY:
    raise RuntimeError("OPENAI_API_KEY 미설정 상태임. 환경변수 설정 필요함.")

client = OpenAI(api_key=OPENAI_API_KEY)

# =========================================================
# 0-1) 상수
# =========================================================
ALLOWED_CLASSES = [
    "바이오","농림수산식품","보건의료","기계","재료","화학공학",
    "전기전자","에너지자원","원자력","환경","건설교통"
]

TEXT_MODEL_FIXED = "gpt-4o-mini" 
IMAGE_MODEL_FIXED = "dall-e-3"    

ICON_SIZE = "1024x1024"
DOWNSCALE_ICON_TO = 512

MAX_PATENT_CHARS = 120_000

OVERVIEW_MIN_CHARS = 50
OVERVIEW_MAX_CHARS = 60

CORE8_MIN_CHARS = 50
CORE8_MAX_CHARS = 60

REP_IMAGE_DPI = 200

TEMPLATE_PPTX_URL = os.environ.get(
    "TEMPLATE_PPTX_URL",
    "https://docs.google.com/presentation/d/12cBPJQlYqf7G44TQkRIRW0r_sBaVcp9I/edit?usp=drive_link&ouid=107110197463763879850&rtpof=true&sd=true"
).strip()

# =========================================================
# 강조 마커 + 색상
# =========================================================
HMARK_L = "⟦"
HMARK_R = "⟧"

COLOR_OLD = RGBColor(0, 176, 80)
COLOR_NEW = RGBColor(68, 114, 196)

_TAGGED_MARK_RE = re.compile(r"⟦([GB]):([^⟦⟧\n]{1,200})⟧")
_MARK_PAIR_RE = re.compile(r"⟦[^⟦⟧\n]{1,200}⟧")

_BAD_EMPH_TOKENS = {
    "기존","본","발명","기술","방법","시스템","구성","제공","가능","적용","요인","한정","리스크","우려",
    "세트","프라이머","검출","진단","구분","활용","관련","사항","존재","내포","가짐","있음",
    "과거","현재","기반","활용해","활용한","활용으로","통해"
}

# =========================================================
# 공통 유틸
# =========================================================
def _shorten_err(e: Exception, maxlen=320) -> str:
    s = str(e).replace("\n", " ").strip()
    return s[:maxlen] + ("..." if len(s) > maxlen else "")

def _norm(x):
    return (x or "").strip()

def extract_json_str(text: str) -> str:
    if not text:
        return ""
    s = text.find("{")
    e = text.rfind("}")
    if s == -1 or e == -1 or e <= s:
        return ""
    return text[s:e+1]

_BAD_END_RE = re.compile(r"(함|됨|임|다|요|니다|습니다|한다|된다|했다|하였다|합니다|입니다)\s*$")

def ends_with_bad_style(s: str) -> bool:
    return bool(_BAD_END_RE.search((s or "").strip()))

def count_mark_pairs_any(s: str) -> int:
    return len(_MARK_PAIR_RE.findall((s or "")))

def count_mark_pairs_tagged(s: str) -> int:
    return len(_TAGGED_MARK_RE.findall((s or "")))

def strip_markers(s: str) -> str:
    s = (s or "")
    s = _TAGGED_MARK_RE.sub(r"\2", s)
    s = s.replace(HMARK_L, "").replace(HMARK_R, "")
    return s

def _len_ws(s: str) -> int:
    return len((s or "").strip())

def _len_ws_no_mark(s: str) -> int:
    return len(strip_markers((s or "").strip()))

def looks_broken_tail(s: str) -> bool:
    t = (strip_markers(s) or "").strip()
    return bool(re.search(
        r"(적용\s*한|운용\s*리스|운용\s*리|적합성\s*적용\s*한|요인\s*적용|한정\s*요인\s*요인)\s*$",
        t
    ))

def clean_korean_title_only(s: str) -> str:
    s = (s or "").strip()
    if not s:
        return ""
    s = re.sub(r"\{[^{}]*\}", "", s).strip()
    s = re.sub(r"\(([A-Za-z0-9 ,.\-_/&:;'\"]+)\)", "", s).strip()
    s = re.sub(r"\[([A-Za-z0-9 ,.\-_/&:;'\"]+)\]", "", s).strip()
    s = re.sub(r"\s{2,}", " ", s).strip()
    return s

def force_noun_ending_minimal(s: str) -> str:
    s = (s or "").strip()
    if not s:
        return s
    s = re.sub(r"[\.!\?]+$", "", s).strip()

    repl = [
        ("했습니다", " 경과"),
        ("하였다", " 경과"),
        ("했다", " 경과"),
        ("합니다", " 기준"),
        ("입니다", " 구조"),
        ("됩니다", " 가능성"),
        ("된다", " 가능성"),
        ("한다", " 수행"),
        ("습니다", " 수준"),
        ("니다", " 수준"),
        ("요", " 요인"),
    ]
    for bad, good in repl:
        if s.endswith(bad):
            s = s[:-len(bad)].rstrip()
            s = re.sub(r"[\.!\?]+$", "", s).strip()
            s = (s + good).strip()
            break

    if s.endswith("함") or s.endswith("됨") or s.endswith("임"):
        s = s[:-1].rstrip()
        s = (s + " 요인").strip()

    return s

def _clip_to_max_word_boundary(s: str, n: int) -> str:
    s = (s or "").strip()
    if len(s) <= n:
        return s
    cut = s[:n].rstrip()
    window = cut[-12:]
    back = max(window.rfind(" "), window.rfind(","), window.rfind(")"), window.rfind("]"))
    if back != -1:
        cut = cut[:len(cut) - (len(window) - back)].rstrip()
    if len(cut) < max(10, n - 15):
        cut = s[:n].rstrip()
    return cut

def clip_only(s: str, mx: int) -> str:
    return _clip_to_max_word_boundary((s or "").strip(), mx)

def force_range_preserve_markers_cliponly(s: str, mx: int) -> str:
    s = (s or "").strip()
    if not s:
        return s

    def safe_chop_tail(x: str) -> str:
        x = x.rstrip()
        if not x:
            return x
        if x.endswith("⟧"):
            idx = x.rfind("⟦")
            if idx > 0:
                head = x[:idx].rstrip()
                if head:
                    head = head[:-1].rstrip()
                    return head + x[idx:]
        return x[:-1].rstrip()

    while _len_ws_no_mark(s) > mx and len(s) > 0:
        s = safe_chop_tail(s)

    s = re.sub(r"\s{2,}", " ", s).strip()
    return s

def enforce_one_mark_pair_plain(s: str) -> str:
    s = (s or "").strip()
    if not s:
        return s

    s = re.sub(r"\*\*([^*\n]{1,200})\*\*", r"\1", s)
    s = re.sub(r"(?<!\*)\*(?!\*)([^*\n]+)(?<!\*)\*(?!\*)", r"\1", s)

    if _TAGGED_MARK_RE.search(s):
        matches = _TAGGED_MARK_RE.findall(s)
        keep_inner = matches[0][1]
        base = _TAGGED_MARK_RE.sub(r"\2", s)
        if keep_inner and keep_inner in base:
            base = re.sub(re.escape(keep_inner), f"{HMARK_L}{keep_inner}{HMARK_R}", base, count=1)
        else:
            base = f"{HMARK_L}{keep_inner}{HMARK_R} " + base
        return re.sub(r"\s{2,}", " ", base).strip()

    raw_pairs = re.findall(r"⟦([^⟦⟧\n]{1,200})⟧", s)
    if len(raw_pairs) >= 1:
        keep = raw_pairs[0]
        base = re.sub(r"⟦([^⟦⟧\n]{1,200})⟧", r"\1", s)
        if keep and keep in base:
            base = re.sub(re.escape(keep), f"{HMARK_L}{keep}{HMARK_R}", base, count=1)
        else:
            base = f"{HMARK_L}{keep}{HMARK_R} " + base
        return re.sub(r"\s{2,}", " ", base).strip()

    tokens = [
        t for t in re.split(r"\s+", s)
        if t and not re.search(r"\d", t) and t not in _BAD_EMPH_TOKENS
    ]
    cand = ""
    for t in tokens:
        if 6 <= len(t) <= 16:
            cand = t
            break
    if not cand and tokens:
        cand = tokens[0][:12]

    if cand and cand in s:
        s = re.sub(re.escape(cand), f"{HMARK_L}{cand}{HMARK_R}", s, count=1)
    else:
        head = s[:12].strip()
        if head:
            s = f"{HMARK_L}{head}{HMARK_R}" + s[len(head):]

    return re.sub(r"\s{2,}", " ", s).strip()

def tag_mark_pair(s: str, tag: str) -> str:
    s = (s or "").strip()
    if not s:
        return s
    tag = "G" if tag.upper() == "G" else "B"

    if _TAGGED_MARK_RE.search(s):
        allm = _TAGGED_MARK_RE.findall(s)
        if len(allm) > 1:
            first_inner = allm[0][1]
            base = _TAGGED_MARK_RE.sub(r"\2", s)
            if first_inner and first_inner in base:
                base = re.sub(re.escape(first_inner), f"⟦{tag}:{first_inner}⟧", base, count=1)
            else:
                base = f"⟦{tag}:{first_inner}⟧ " + base
            return re.sub(r"\s{2,}", " ", base).strip()
        m = _TAGGED_MARK_RE.search(s)
        if m:
            inner = m.group(2)
            return _TAGGED_MARK_RE.sub(lambda mm: f"⟦{tag}:{inner}⟧", s, count=1)
        return s

    m = re.search(r"⟦([^⟦⟧\n]{1,200})⟧", s)
    if not m:
        return s
    inner = m.group(1)
    s = s.replace(f"⟦{inner}⟧", f"⟦{tag}:{inner}⟧", 1)
    return s

def hard_fix_core8_minimal(s: str, tag: str, mx: int) -> str:
    s = force_noun_ending_minimal(s)
    s = enforce_one_mark_pair_plain(s)
    s = tag_mark_pair(s, tag)
    s = force_range_preserve_markers_cliponly(s, mx)
    s = enforce_one_mark_pair_plain(s)
    s = tag_mark_pair(s, tag)
    s = force_range_preserve_markers_cliponly(s, mx)
    return s

# =========================================================
# 1) 업로드/URL 파일 처리
# =========================================================
def convert_google_slides_to_pptx_export_url(url: str) -> str:
    if not url:
        return url
    m = re.search(r"/presentation/d/([a-zA-Z0-9\-_]+)", url)
    if not m:
        return url
    file_id = m.group(1)
    return f"https://docs.google.com/presentation/d/{file_id}/export/pptx"

def download_file_to_temp(url: str, suffix: str, timeout=60) -> str:
    if not url:
        return ""

    if "docs.google.com/presentation/d/" in url and suffix.lower() == ".pptx":
        url = convert_google_slides_to_pptx_export_url(url)

    cache_key = hashlib.md5(f"{url}|{suffix}".encode("utf-8")).hexdigest()
    cache_path = os.path.join(tempfile.gettempdir(), f"{cache_key}{suffix}")

    if os.path.exists(cache_path) and os.path.getsize(cache_path) > 0:
        return cache_path

    headers = {"User-Agent": "Mozilla/5.0", "Accept": "*/*"}
    with requests.Session() as sess:
        r = sess.get(url, headers=headers, stream=True, timeout=timeout, allow_redirects=True)
        r.raise_for_status()

        content_type = (r.headers.get("Content-Type") or "").lower()
        if "text/html" in content_type and suffix.lower() == ".pptx":
            raise RuntimeError(f"PPTX 다운로드 실패: HTML 응답 수신. url={url}")

        with open(cache_path, "wb") as f:
            for chunk in r.iter_content(chunk_size=1024 * 1024):
                if chunk:
                    f.write(chunk)
    return cache_path

def get_uploaded_path(file_obj, suffix):
    if file_obj is None:
        return None
    if isinstance(file_obj, str) and os.path.exists(file_obj):
        return file_obj
    if hasattr(file_obj, "read"):
        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
            tmp.write(file_obj.read())
            return tmp.name
    return None

# =========================================================
# 2) 텍스트 추출 (PDF & DOCX)
# =========================================================
def extract_pdf_text(pdf_path: str, max_pages: int = 25) -> str:
    doc = fitz.open(pdf_path)
    texts = []
    for i in range(min(len(doc), max_pages)):
        page = doc[i]
        texts.append(page.get_text("text") or "")
    return "\n".join(texts).strip()

def extract_docx_text(docx_path: str) -> str:
    doc = docx.Document(docx_path)
    texts = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n".join(texts).strip()

# =========================================================
# 2-1) 특허 메타데이터 파싱
# =========================================================
def parse_patent_metadata(patent_text: str) -> dict:
    t = patent_text or ""
    t = t.replace("\r", "\n")
    t = re.sub(r"[ \t]+", " ", t)
    t = re.sub(r"\n{2,}", "\n", t)

    head = t[:12000]

    def clean_title(s: str) -> str:
        s = (s or "").strip()
        s = re.sub(r"^[】\]\)\s]+", "", s).strip()
        s = re.split(r"요\s*약|【요약】|【기술분야】|【발명의 배경이 되는 기술】|【배경기술】", s, maxsplit=1)[0].strip()
        s = re.split(r"【[^】]+】", s)[0].strip()
        s = re.sub(r"^\[\d+\]\s*", "", s).strip()
        s = re.sub(r"^【\d+】\s*", "", s).strip()
        s = clean_korean_title_only(s)
        s = re.sub(r"\s{2,}", " ", s).strip()
        return s[:200].strip()

    def find_title() -> str:
        patterns = [
            r"발명의 명칭\s*[:：]?\s*(.+?)(?=요\s*약|【요약】|【기술분야】|【발명의 배경이 되는 기술】|【배경기술】|\Z)",
            r"【발명의 명칭】\s*(.+?)(?=【기술분야】|【발명의 배경이 되는 기술】|【배경기술】|\Z)",
            r"발명의 국문명칭\s*[:：]?\s*(.+)",
            r"\(54\)\s*발명의 명칭\s*(.+?)(?=\(57\)|요\s*약|\Z)",
        ]
        for pat in patterns:
            m = re.search(pat, head, flags=re.DOTALL)
            if m:
                title = clean_title(m.group(1))
                if 2 <= len(title) <= 150:
                    return title

        lines = [x.strip() for x in head.split("\n") if x.strip()]
        for line in lines[:40]:
            if "발명의 명칭" in line or "발명의 국문명칭" in line:
                continue
            cand = clean_title(line)
            if 4 <= len(cand) <= 120:
                return cand
        return ""

    def find_after_label(label: str) -> str:
        patterns = [
            rf"{label}\s*[:：]?\s*([^\n]+)",
            rf"{label}\s*\n\s*([^\n]+)",
            rf"【{label}】\s*([^\n]+)",
        ]
        for pat in patterns:
            m = re.search(pat, t)
            if m:
                return (m.group(1) or "").strip()
        return ""

    def pick_number_like(s: str) -> str:
        s = (s or "").strip()
        if not s:
            return ""
        patterns = [
            r"(10-\d{4}-\d{6,7})",
            r"(10-\d{7})",
            r"(\d{2,4}-\d{2,4}-\d{2,})",
            r"(\d{8,})",
        ]
        for pat in patterns:
            m = re.search(pat, s)
            if m:
                return m.group(1).replace("–", "-")
        return ""

    def pick_date_like(s: str) -> str:
        s = (s or "").strip()
        if not s:
            return ""
        nums = re.findall(r"\d+", s)
        if len(nums) >= 3:
            y, m, d = nums[0], nums[1], nums[2]
            try:
                return f"{int(y):04d}.{int(m):02d}.{int(d):02d}"
            except Exception:
                return ""
        return ""

    invention_title = clean_korean_title_only(find_title())

    app_no = pick_number_like(find_after_label("출원번호"))
    reg_no = pick_number_like(find_after_label("등록번호"))
    app_date = pick_date_like(find_after_label("출원일자") or find_after_label("출원일"))
    reg_date = pick_date_like(find_after_label("등록일자") or find_after_label("등록일"))

    if app_no and reg_no:
        app_reg_no_display = f"{app_no}({reg_no})"
    elif app_no:
        app_reg_no_display = app_no
    elif reg_no:
        app_reg_no_display = reg_no
    else:
        app_reg_no_display = ""

    if app_date and reg_date:
        app_reg_date_display = f"{app_date}({reg_date})"
    elif app_date:
        app_reg_date_display = app_date
    elif reg_date:
        app_reg_date_display = reg_date
    else:
        app_reg_date_display = ""

    return {
        "invention_title": invention_title,
        "app_no": app_no,
        "reg_no": reg_no,
        "app_date": app_date,
        "reg_date": reg_date,
        "app_reg_no_display": app_reg_no_display,
        "app_reg_date_display": app_reg_date_display,
    }

# =========================================================
# 3) OpenAI 규칙
# =========================================================
def _supports_response_format() -> bool:
    try:
        sig = inspect.signature(client.chat.completions.create) 
        return "response_format" in sig.parameters
    except Exception:
        return False

def _repair_prompt_for_targets(targets: dict, stage_name: str) -> str:
    return f"""
아래 JSON 문장들을 규칙에 맞게 재작성하라.
중요: 길이 미달/초과를 "절단/패딩/말미 꼬리 덧붙이기"로 해결 금지.
반드시 "문장 내부의 명사구 보강(수식어/관형어/명사구 추가)"으로 길이를 맞추어라.
추가되는 정보는 반드시 원문 의미 범위 안에서만 확장하고, 사실을 새로 발명하지 말 것.

[공통 규칙]
- 각 문장 1문장 유지, 줄바꿈 금지
- 문장 종결은 반드시 명사(명사구)로 종료
- 금지 종결: "~함", "~됨", "~임"
- 서술형 어미 사용 금지
- “→” 금지, ":" ";" 중심 나열 금지
- 단일 별표(*...*) 금지, **...** 금지
- 길이 조정은 "문장 내부"에서만 수행

[길이 규칙]
- overview: 공백 포함 {OVERVIEW_MIN_CHARS}~{OVERVIEW_MAX_CHARS}자
- 핵심8: 표식 길이 제외 {CORE8_MIN_CHARS}~{CORE8_MAX_CHARS}자

[핵심8 강조 규칙]
- 각 문장에 ⟦강조구절⟧ 1개 정확히 1개 포함
- ⟦...⟧ 내부 구절은 문장 내에서 의미가 가장 핵심인 명사구로 선정
- 길이 조정 과정에서도 ⟦...⟧ 구절 의미 변경 금지
- ⟦...⟧ 내부에는 키워드 나열 금지
- ⟦...⟧ 는 반드시 하나의 명사구만 포함

[입력 JSON]
{json.dumps(targets, ensure_ascii=False, indent=2)}

[출력 스키마]
{{
  "overview": [{{"t":"string","n":0}},{{"t":"string","n":0}},{{"t":"string","n":0}}],
  "diff_old": [{{"t":"string","n":0}},{{"t":"string","n":0}}],
  "diff_new": [{{"t":"string","n":0}},{{"t":"string","n":0}}],
  "limits_old": [{{"t":"string","n":0}},{{"t":"string","n":0}}],
  "advantages_new": [{{"t":"string","n":0}},{{"t":"string","n":0}}]
}}
""".strip()

def _pick_text_list(obj_list, expected_len):
    if not isinstance(obj_list, list) or len(obj_list) != expected_len:
        return None
    out = []
    for it in obj_list:
        if isinstance(it, dict) and isinstance(it.get("t"), str):
            out.append(_norm(it["t"]))
        else:
            return None
    return out

def repair_targets_via_llm(payload: dict, log: list, stage: str) -> dict:
    targets = {
        "overview": payload.get("overview", ["","",""]),
        "diff_old": payload.get("diff_old", ["",""]),
        "diff_new": payload.get("diff_new", ["",""]),
        "limits_old": payload.get("limits_old", ["",""]),
        "advantages_new": payload.get("advantages_new", ["",""]),
    }
    prompt = _repair_prompt_for_targets(targets, stage)

    try:
        r = client.chat.completions.create(
            model=TEXT_MODEL_FIXED,
            temperature=0.0,
            messages=[{"role": "user", "content": prompt}],
            max_tokens=1200,
        )
        js = extract_json_str(r.choices[0].message.content or "")
        if not js:
            log.append(f"- {stage} 실패: JSON 추출 실패")
            return payload
        fixed = json.loads(js)

        for k, n_expected in [("overview",3),("diff_old",2),("diff_new",2),("limits_old",2),("advantages_new",2)]:
            picked = _pick_text_list(fixed.get(k), n_expected)
            if picked is not None:
                payload[k] = picked

        log.append(f"- {stage} 적용: 문장 내부 명사구 보강 리페어 수행")
        return payload
    except Exception as e:
        log.append(f"- {stage} 실패: {type(e).__name__}: {_shorten_err(e)}")
        return payload

def _check_len_rules(payload: dict) -> list:
    issues = []
    if not isinstance(payload, dict):
        return ["payload dict 아님"]

    ov = payload.get("overview", [])
    if not isinstance(ov, list) or len(ov) != 3:
        issues.append("overview 형식 오류")
    else:
        bad = sum(1 for s in ov if _len_ws(s) < OVERVIEW_MIN_CHARS or _len_ws(s) > OVERVIEW_MAX_CHARS)
        if bad:
            issues.append(f"overview 길이 위반 {bad}개")

    for k in ["diff_old","diff_new","limits_old","advantages_new"]:
        arr = payload.get(k, [])
        if not isinstance(arr, list) or len(arr) != 2:
            issues.append(f"{k} 형식 오류")
            continue
        bad = sum(1 for s in arr if _len_ws_no_mark(s) < CORE8_MIN_CHARS or _len_ws_no_mark(s) > CORE8_MAX_CHARS)
        if bad:
            issues.append(f"{k} 길이 위반 {bad}개")
        mw = sum(1 for s in arr if count_mark_pairs_any(s) != 1)
        if mw:
            issues.append(f"{k} 마커 위반 {mw}개")
    return issues

def enforce_len_by_repair_loop(payload: dict, log: list, max_rounds: int = 1) -> dict:
    for r in range(1, max_rounds + 1):
        issues = _check_len_rules(payload)
        if not issues:
            log.append(f"- 길이 규칙 만족 확인: round={r-1}")
            return payload
        log.append(f"- 길이 규칙 미달 감지: round={r} / issues={issues}")
        payload = repair_targets_via_llm(payload, log, stage=f"리페어 라운드 {r}")

    log.append("- 길이 규칙 최종 미만족: max_rounds 도달")
    return payload

def validate_and_normalize_payload_minimal(payload: dict):
    issues = []
    if not isinstance(payload, dict):
        return {}, ["payload가 dict가 아님"]

    out = {
        "tech_name": _norm(payload.get("tech_name","")),
        "tech_class": _norm(payload.get("tech_class","")),
        "tech_summary": _norm(payload.get("tech_summary","")),
        "applications": payload.get("applications", []),
        "overview": payload.get("overview", []),
        "diff_old": payload.get("diff_old", []),
        "diff_old_extra": payload.get("diff_old_extra", ["",""]),
        "diff_new": payload.get("diff_new", []),
        "diff_new_extra": payload.get("diff_new_extra", ["",""]),
        "limits_old": payload.get("limits_old", []),
        "advantages_new": payload.get("advantages_new", []),
    }

    if out["tech_class"] not in ALLOWED_CLASSES:
        issues.append(f"tech_class 허용목록 위반: {out['tech_class']}")
        out["tech_class"] = ""

    apps = out["applications"] if isinstance(out["applications"], list) else []
    while len(apps) < 3:
        apps.append({"title":"","desc":""})
        issues.append("applications 항목 3개 미만 → 빈 항목 보정")
    apps = apps[:3]
    out["applications"] = [{
        "title": _norm(a.get("title","")) if isinstance(a, dict) else "",
        "desc": _norm(a.get("desc","")) if isinstance(a, dict) else ""
    } for a in apps]

    def ensure_list(val, n):
        val = val if isinstance(val, list) else []
        val = [_norm(x) for x in val]
        while len(val) < n:
            val.append("")
        return val[:n]

    out["overview"] = ensure_list(out["overview"], 3)
    out["diff_old"] = ensure_list(out["diff_old"], 2)
    out["diff_old_extra"] = ensure_list(out["diff_old_extra"], 2)
    out["diff_new"] = ensure_list(out["diff_new"], 2)
    out["diff_new_extra"] = ensure_list(out["diff_new_extra"], 2)
    out["limits_old"] = ensure_list(out["limits_old"], 2)
    out["advantages_new"] = ensure_list(out["advantages_new"], 2)

    out["overview"] = [force_noun_ending_minimal(s) for s in out["overview"]]
    out["overview"] = [clip_only(s, OVERVIEW_MAX_CHARS) for s in out["overview"]]

    out["diff_old"] = [hard_fix_core8_minimal(s, "G", CORE8_MAX_CHARS) for s in out["diff_old"]]
    out["limits_old"] = [hard_fix_core8_minimal(s, "G", CORE8_MAX_CHARS) for s in out["limits_old"]]
    out["diff_new"] = [hard_fix_core8_minimal(s, "B", CORE8_MAX_CHARS) for s in out["diff_new"]]
    out["advantages_new"] = [hard_fix_core8_minimal(s, "B", CORE8_MAX_CHARS) for s in out["advantages_new"]]

    return out, issues

def call_openai_json(patent_text: str, temperature: float, log: list, meta: dict = None) -> dict:
    patent_text = (patent_text or "")[:MAX_PATENT_CHARS]
    model_name = TEXT_MODEL_FIXED
    invention_title = clean_korean_title_only((meta or {}).get("invention_title", ""))

    log.append("- 안내: 호환 모드(JSON 강제+추출)로 진행")

    compat_prompt = f"""
아래 규칙을 모두 만족하여 JSON만 출력하라. JSON 외 텍스트 출력 금지.

[중요 지시]
- tech_name은 단순히 발명의 명칭을 반복하거나 축약한 표현이 아니라,
  수요기업이 한눈에 보고 기술의 기능과 적용 맥락을 판단하기 쉬운 "쉬운 기술제목"으로 작성
- tech_name은 마케팅 문구처럼 과장하지 말고, 실제 기술 기능·용도·적용 대상을 반영한 직관적 제목으로 작성
- tech_name은 가능한 한 한글 중심으로 작성
- 참고 발명의 명칭: "{invention_title}"

[규칙]
- tech_class는 다음 중 1개만 허용: {" / ".join(ALLOWED_CLASSES)}
- applications는 반드시 3개
- overview는 반드시 3개이며 각 문장 공백 포함 {OVERVIEW_MIN_CHARS}~{OVERVIEW_MAX_CHARS}자
- diff_old, diff_new는 각각 2개
- diff_old_extra, diff_new_extra는 각각 2개
- limits_old, advantages_new는 각각 2개
- 모든 문장 명사형 종결
- "함/됨/임" 종결 금지
- 서술형 어미 사용 금지
- “→” 기호 절대 사용 금지
- 단일 별표 강조(*...*) 금지, **...** 사용 금지
- diff_old/diff_new/limits_old/advantages_new 8문장은 (표식 길이 제외) {CORE8_MIN_CHARS}~{CORE8_MAX_CHARS}자
- 위 8문장 각각에 ⟦강조구절⟧ 1개를 정확히 1개만 포함

[강조 규칙]
- ⟦...⟧ 내부 구절은 문장에서 가장 핵심적인 기술 개념 명사구 선택
- 기능, 구조, 핵심 기술원리, 차별요소 중 하나를 강조
- "기존기술/본 기술/방법/기술/시스템" 같은 일반 단어 강조 금지
- 키워드 나열 금지
- 반드시 하나의 의미 있는 명사구만 포함

- 길이 조정은 절단/패딩 금지, 문장 내부 명사구 보강으로만 수행

[JSON 스키마]
{{
  "tech_name": "string",
  "tech_class": "string",
  "tech_summary": "string",
  "applications": [
    {{"title":"string","desc":"string"}},
    {{"title":"string","desc":"string"}},
    {{"title":"string","desc":"string"}}
  ],
  "overview": ["string","string","string"],
  "diff_old": ["string","string"],
  "diff_old_extra": ["string","string"],
  "diff_new": ["string","string"],
  "diff_new_extra": ["string","string"],
  "limits_old": ["string","string"],
  "advantages_new": ["string","string"]
}}

[특허 명세서 원문]
{patent_text}
""".strip()

    resp = client.chat.completions.create(
        model=model_name,
        temperature=float(temperature),
        messages=[{"role": "user", "content": compat_prompt}],
        max_tokens=2400,
    )
    raw = resp.choices[0].message.content or ""
    js = extract_json_str(raw)
    if not js:
        raise RuntimeError("호환 모드 JSON 추출 실패")

    payload = json.loads(js)
    payload = enforce_len_by_repair_loop(payload, log=log, max_rounds=1)
    payload, issues = validate_and_normalize_payload_minimal(payload)
    log.append("- JSON 경고 없음" if not issues else f"- JSON 경고: {issues}")
    return payload

# =========================================================
# 4) JSON → 가~사 텍스트
# =========================================================
def payload_to_readable_text(p: dict) -> str:
    apps = p["applications"]
    meta = p.get("meta", {}) if isinstance(p.get("meta", {}), dict) else {}

    def S(x): return strip_markers(x)

    lines = []
    lines.append("0. 문서 기본정보")
    lines.append(f"- 발명의 명칭: {meta.get('invention_title','')}")
    lines.append(f"- 출원(등록)번호: {meta.get('app_reg_no_display','')}")
    lines.append(f"- 출원(등록)일자: {meta.get('app_reg_date_display','')}")
    lines.append(f"- 출원번호: {meta.get('app_no','')}")
    lines.append(f"- 출원일자: {meta.get('app_date','')}")
    lines.append(f"- 등록번호: {meta.get('reg_no','')}")
    lines.append(f"- 등록일자: {meta.get('reg_date','')}")
    lines.append("")

    lines.append("가. 기술명"); lines.append(p["tech_name"]); lines.append("")
    lines.append("나. 기술분류"); lines.append(p["tech_class"]); lines.append("")
    lines.append("다. 기술 요약"); lines.append(p["tech_summary"]); lines.append("")
    lines.append("라. 적용분야 및 제품")
    for idx, a in enumerate(apps, start=1):
        lines.append(f"{idx}) {a['title']} - {a['desc']}")
    lines.append("")
    lines.append("마. 기술 개요")
    for i, s in enumerate(p["overview"], start=1):
        lines.append(f"{i}) {s} ({_len_ws(s)}자)")
    lines.append("")
    lines.append("바. 기술 차별성")
    lines.append("[기존기술]")
    lines.append(S(p["diff_old"][0])); lines.append(S(p["diff_old"][1]))
    if p["diff_old_extra"][0] or p["diff_old_extra"][1]:
        lines.append(p["diff_old_extra"][0]); lines.append(p["diff_old_extra"][1])
    lines.append("")
    lines.append("[본 기술]")
    lines.append(S(p["diff_new"][0])); lines.append(S(p["diff_new"][1]))
    if p["diff_new_extra"][0] or p["diff_new_extra"][1]:
        lines.append(p["diff_new_extra"][0]); lines.append(p["diff_new_extra"][1])
    lines.append("")
    lines.append("사. 기술적 한계 및 우위")
    lines.append("[기존기술(한계)]")
    lines.append(f"1) {S(p['limits_old'][0])}")
    lines.append(f"2) {S(p['limits_old'][1])}")
    lines.append("")
    lines.append("[본 기술(우위)]")
    lines.append(f"1) {S(p['advantages_new'][0])}")
    lines.append(f"2) {S(p['advantages_new'][1])}")
    return "\n".join(lines).strip()

# =========================================================
# 5) PPT 플레이스홀더 치환
# =========================================================
def _iter_shapes_recursive(shapes):
    for sh in shapes:
        yield sh
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            for sub in _iter_shapes_recursive(sh.shapes):
                yield sub

def _replace_in_paragraph_runs(paragraph, mapping, counts: dict):
    runs = list(paragraph.runs)
    if not runs:
        return

    changed = True
    while changed:
        changed = False
        full = "".join(r.text for r in runs)
        if "{{" not in full:
            break

        for key, val in mapping.items():
            if val is None:
                continue
            val = str(val)

            start = full.find(key)
            if start == -1:
                continue
            end = start + len(key)

            pos = 0
            span = []
            for idx, r in enumerate(runs):
                r_start = pos
                r_end = pos + len(r.text)
                if r_end > start and r_start < end:
                    span.append((idx, max(start, r_start) - r_start, min(end, r_end) - r_start))
                pos = r_end

            if not span:
                continue

            first_i, first_s, _ = span[0]
            last_i, _, last_e = span[-1]

            prefix = runs[first_i].text[:first_s]
            suffix = runs[last_i].text[last_e:]

            runs[first_i].text = prefix + val + suffix
            for j in range(first_i + 1, last_i + 1):
                runs[j].text = ""

            counts[key] = counts.get(key, 0) + 1
            changed = True
            break

def _replace_in_text_frame(text_frame, mapping, counts: dict):
    for paragraph in text_frame.paragraphs:
        _replace_in_paragraph_runs(paragraph, mapping, counts)

def replace_placeholders_everywhere(prs, mapping):
    counts = {}
    for slide in prs.slides:
        for shape in _iter_shapes_recursive(slide.shapes):
            if getattr(shape, "has_text_frame", False):
                _replace_in_text_frame(shape.text_frame, mapping, counts)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        _replace_in_text_frame(cell.text_frame, mapping, counts)
    return counts

# =========================================================
# 6) 남은 플레이스홀더 스캔
# =========================================================
_PLACEHOLDER_RE = re.compile(r"\{\{[^}]+\}\}")

def scan_remaining_placeholders_recursive(prs):
    remain = set()

    def scan_text(text: str):
        if not text:
            return
        for m in _PLACEHOLDER_RE.findall(text):
            remain.add(m)

    for slide in prs.slides:
        for shape in _iter_shapes_recursive(slide.shapes):
            if getattr(shape, "has_text_frame", False):
                scan_text(shape.text or "")
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        scan_text(cell.text or "")
    return sorted(remain)

def check_placeholders_status(mapping: dict, replace_counts: dict, remaining_placeholders: list, log: list, must_have_keys=None):
    must_have_keys = must_have_keys or []

    log.append("[치환 리포트] placeholder별 치환 횟수")
    for k in mapping.keys():
        c = int(replace_counts.get(k, 0))
        log.append(f"- {k}: {c}회")

    if remaining_placeholders:
        log.append(f"[경고] 남은 플레이스홀더 발견: {remaining_placeholders}")
        for k in must_have_keys:
            if k in remaining_placeholders:
                log.append(f"[강경고] 필수 메타 플레이스홀더가 남아있음: {k}")
    else:
        log.append("[치환 리포트] 남은 플레이스홀더 없음")

# =========================================================
# 7) PPT 도형 검색/제거
# =========================================================
def find_placeholder_shape(prs, placeholder):
    for slide in prs.slides:
        for shape in list(slide.shapes):
            if getattr(shape, "has_text_frame", False):
                if placeholder in (shape.text or ""):
                    return slide, shape
    return None, None

def remove_shape(slide, shape):
    slide.shapes._spTree.remove(shape._element)

# =========================================================
# 8) OpenAI 아이콘 생성 + PPT 삽입
# =========================================================
ICON_STYLE_PROMPT = (
    "single simple outline icon, bold clean black stroke line icon, isolated on pure white background, "
    "one object only, centered composition, very simplified silhouette, minimal vector look, "
    "no shadow, no gradient, no text, no label, no number, no watermark"
)

def generate_icon_png_stream(prompt: str) -> io.BytesIO:
    img = client.images.generate(
        model=IMAGE_MODEL_FIXED,
        prompt=prompt,
        size=ICON_SIZE,
    )

    if not getattr(img, "data", None):
        raise RuntimeError("이미지 응답 데이터 없음")

    first = img.data[0]
    b64 = getattr(first, "b64_json", None)
    url = getattr(first, "url", None)
    
    if b64:
        png_bytes = base64.b64decode(b64)
    elif url:
        resp = requests.get(url)
        png_bytes = resp.content
    else:
        raise RuntimeError("이미지 응답에 b64_json 또는 url이 없음")

    pil = Image.open(io.BytesIO(png_bytes)).convert("RGBA")

    if DOWNSCALE_ICON_TO and isinstance(DOWNSCALE_ICON_TO, int):
        pil = pil.resize((DOWNSCALE_ICON_TO, DOWNSCALE_ICON_TO), Image.LANCZOS)

    out = io.BytesIO()
    pil.save(out, format="PNG")
    out.seek(0)
    return out

def hash_png_stream(png_stream: io.BytesIO) -> str:
    return hashlib.sha256(png_stream.getvalue()).hexdigest()

def _make_one_icon(i: int, payload: dict):
    title = payload["applications"][i]["title"]
    desc = payload["applications"][i]["desc"]
    p = f"{ICON_STYLE_PROMPT}. Subject: {title}. Context: {desc}. Make it distinct from other icons."
    stream = generate_icon_png_stream(p)
    h = hash_png_stream(stream)
    return i, stream, h

def generate_three_icons(payload: dict, log: list):
    icons = [None, None, None]
    hashes = [None, None, None]

    with ThreadPoolExecutor(max_workers=3) as ex:
        futures = [ex.submit(_make_one_icon, i, payload) for i in range(3)]
        for fut in as_completed(futures):
            i, stream, h = fut.result()
            icons[i] = stream
            hashes[i] = h

    seen = set()
    for i in range(3):
        if hashes[i] in seen:
            log.append(f"- 아이콘 중복 감지(i={i+1}) → 1회 재생성")
            title = payload["applications"][i]["title"]
            desc = payload["applications"][i]["desc"]
            p2 = (
                f"{ICON_STYLE_PROMPT}. Subject: {title}. Context: {desc}. "
                "Ensure a clearly different silhouette and object category than the previous icons."
            )
            icons[i] = generate_icon_png_stream(p2)
            hashes[i] = hash_png_stream(icons[i])
        seen.add(hashes[i])

    return icons

def replace_icons_with_streams(prs, icon_streams, log):
    placeholders = ["{{아이콘 1}}", "{{아이콘 2}}", "{{아이콘 3}}"]
    for ph, stream in zip(placeholders, icon_streams):
        slide, shape = find_placeholder_shape(prs, ph)
        if not shape:
            log.append(f"- {ph}: 템플릿에서 플레이스홀더 미발견")
            continue
        left, top = shape.left, shape.top
        remove_shape(slide, shape)
        slide.shapes.add_picture(stream, left, top, width=Cm(1.5), height=Cm(1.5))
        log.append(f"- {ph}: 삽입 완료")

# =========================================================
# 9) 대표이미지 추출
# =========================================================
def _normalize_text_for_search(s: str) -> str:
    s = (s or "").replace("\r", "\n")
    s = re.sub(r"[ \t]+", " ", s)
    s = re.sub(r"\n{2,}", "\n", s)
    return s.strip()

def extract_text_all_pages(doc):
    texts = []
    for i in range(len(doc)):
        try:
            texts.append(doc[i].get_text("text") or "")
        except Exception:
            texts.append("")
    return texts

def trim_header_footer(pil_img, top_ratio=0.01, bottom_ratio=0.01):
    w, h = pil_img.size
    top = int(h * top_ratio)
    bottom = int(h * (1 - bottom_ratio))
    if bottom <= top:
        return pil_img
    return pil_img.crop((0, top, w, bottom))

def remove_footer_band(pil_img, ratio=0.03):
    w, h = pil_img.size
    cut = int(h * (1 - ratio))
    if cut <= 0:
        return pil_img
    return pil_img.crop((0, 0, w, cut))

def autocrop_whitespace(pil_img, bg_color=(255,255,255), tolerance=8):
    img = pil_img.convert("RGB") if pil_img.mode != "RGB" else pil_img
    bg = Image.new("RGB", img.size, bg_color)
    diff = ImageChops.difference(img, bg)
    bbox = diff.point(lambda x: 255 if x > tolerance else 0).getbbox()
    return img.crop(bbox) if bbox else img

def autocrop_main_figure_with_caption(pil_img, bg_threshold=245, min_area_ratio=0.0015, pad=10, caption_gap_ratio=0.18, footer_exclude_ratio=0.90):
    img = pil_img.convert("RGB")
    gray = img.convert("L")
    w, h = gray.size

    mask = gray.point(lambda x: 0 if x >= bg_threshold else 255, mode="L")
    arr = np.array(mask, dtype=np.uint8)
    fg = arr > 0

    if not fg.any():
        return img

    visited = np.zeros_like(fg, dtype=bool)
    components = []

    def neighbors(y, x):
        if y > 0: yield y - 1, x
        if y < h - 1: yield y + 1, x
        if x > 0: yield y, x - 1
        if x < w - 1: yield y, x + 1

    for y in range(h):
        for x in range(w):
            if not fg[y, x] or visited[y, x]:
                continue

            stack = [(y, x)]
            visited[y, x] = True

            min_x = max_x = x
            min_y = max_y = y
            pixel_area = 0

            while stack:
                cy, cx = stack.pop()
                pixel_area += 1

                if cx < min_x: min_x = cx
                if cx > max_x: max_x = cx
                if cy < min_y: min_y = cy
                if cy > max_y: max_y = cy

                for ny, nx in neighbors(cy, cx):
                    if fg[ny, nx] and not visited[ny, nx]:
                        visited[ny, nx] = True
                        stack.append((ny, nx))

            width = max_x - min_x + 1
            height = max_y - min_y + 1
            box_area = width * height
            cx = (min_x + max_x) / 2.0
            cy = (min_y + max_y) / 2.0

            components.append({
                "bbox": (min_x, min_y, max_x + 1, max_y + 1),
                "pixel_area": pixel_area,
                "box_area": box_area,
                "width": width,
                "height": height,
                "cx": cx,
                "cy": cy,
            })

    if not components:
        return img

    min_area = w * h * float(min_area_ratio)
    filtered = [c for c in components if c["pixel_area"] >= min_area]
    if not filtered:
        filtered = sorted(components, key=lambda c: c["pixel_area"], reverse=True)[:5]

    def main_score(c):
        score = (
            c["pixel_area"] * 1.0 +
            c["box_area"] * 0.03 +
            c["width"] * 0.3 +
            c["height"] * 0.3
        )
        if c["cy"] > h * footer_exclude_ratio:
            score *= 0.2
        return score

    main = max(filtered, key=main_score)
    mx0, my0, mx1, my1 = main["bbox"]

    selected = [main]
    max_caption_gap = int(h * caption_gap_ratio)

    for c in filtered:
        if c is main:
            continue

        x0, y0, x1, y1 = c["bbox"]

        if y0 >= h * footer_exclude_ratio:
            continue

        vertical_gap = y0 - my1
        horizontally_related = not (x1 < mx0 - 40 or x0 > mx1 + 40)

        likely_caption = (
            0 <= vertical_gap <= max_caption_gap
            and horizontally_related
            and c["height"] < h * 0.12
            and c["width"] < w * 0.95
        )

        near_body = (
            y1 >= my0 - 20 and y0 <= my1 + 20 and
            x1 >= mx0 - 20 and x0 <= mx1 + 20
        )

        if likely_caption or near_body:
            selected.append(c)

    ux0 = min(c["bbox"][0] for c in selected)
    uy0 = min(c["bbox"][1] for c in selected)
    ux1 = max(c["bbox"][2] for c in selected)
    uy1 = max(c["bbox"][3] for c in selected)

    ux0 = max(0, ux0 - pad)
    uy0 = max(0, uy0 - pad)
    ux1 = min(w, ux1 + pad)
    uy1 = min(h, uy1 + pad)

    cropped = img.crop((ux0, uy0, ux1, uy1))
    cropped = autocrop_whitespace(cropped)
    return cropped

def is_reasonable_figure_image(pil_img, min_w=220, min_h=160):
    img = pil_img.convert("RGB")
    gray = img.convert("L")
    w, h = gray.size

    if w < min_w or h < min_h:
        return False, f"too_small:{w}x{h}"

    arr = np.array(gray, dtype=np.uint8)

    fg = arr < 245
    fg_ratio = float(fg.mean())
    if fg_ratio < 0.01:
        return False, f"fg_too_low:{fg_ratio:.4f}"

    ys, xs = np.where(fg)
    if len(xs) == 0 or len(ys) == 0:
        return False, "no_fg"

    bw = xs.max() - xs.min() + 1
    bh = ys.max() - ys.min() + 1

    if bw < min_w * 0.6 or bh < min_h * 0.6:
        return False, f"bbox_too_small:{bw}x{bh}"

    small = img.resize((max(32, min(96, w // 4)), max(32, min(96, h // 4))))
    colors = small.convert("P", palette=Image.ADAPTIVE, colors=16).getcolors()
    color_count = len(colors) if colors else 0

    if color_count <= 2 and fg_ratio < 0.08:
        return False, f"too_simple_colors:{color_count}/fg={fg_ratio:.4f}"

    bbox_area_ratio = (bw * bh) / float(w * h)
    if bbox_area_ratio < 0.12:
        return False, f"bbox_area_too_small:{bbox_area_ratio:.4f}"

    return True, "ok"

def render_page_to_pil(page, dpi=REP_IMAGE_DPI, clip_rect=None, apply_trim=True, apply_footer_cut=True, main_figure_crop=True):
    mat = fitz.Matrix(dpi / 72, dpi / 72)
    pix = page.get_pixmap(matrix=mat, alpha=False, clip=clip_rect)
    pil = Image.open(io.BytesIO(pix.tobytes("png"))).convert("RGB")

    if apply_trim:
        pil = trim_header_footer(pil, top_ratio=0.01, bottom_ratio=0.01)

    if apply_footer_cut:
        pil = remove_footer_band(pil, ratio=0.03)

    if main_figure_crop:
        pil = autocrop_main_figure_with_caption(pil)
    else:
        pil = autocrop_whitespace(pil)

    return pil

def pil_to_png_stream(pil):
    out = io.BytesIO()
    pil.save(out, format="PNG")
    out.seek(0)
    return out

def extract_representative_from_pdf_registered_style(pdf_path, dpi=REP_IMAGE_DPI):
    doc = fitz.open(pdf_path)
    page = doc[0]

    imgs = page.get_images(full=True)
    best_bytes = None
    best_area = 0

    for img in imgs:
        try:
            xref = img[0]
            info = doc.extract_image(xref)
            b = info.get("image")
            if not b:
                continue
            pil = Image.open(io.BytesIO(b))
            area = pil.size[0] * pil.size[1]
            if area > best_area:
                best_area = area
                best_bytes = b
        except Exception:
            continue

    if best_bytes:
        pil_img = Image.open(io.BytesIO(best_bytes)).convert("RGB")
        pil_img = autocrop_whitespace(pil_img)
        return pil_to_png_stream(pil_img), {
            "strategy": "registered_style_first_page_image",
            "page_index": 0,
            "reason": "등록공보 감지 → 첫 페이지 최대 이미지 추출"
        }

    mat = fitz.Matrix(dpi / 72, dpi / 72)
    pix = page.get_pixmap(matrix=mat, alpha=False)
    pil_img = Image.open(io.BytesIO(pix.tobytes("png"))).convert("RGB")
    pil_img = autocrop_whitespace(pil_img)
    return pil_to_png_stream(pil_img), {
        "strategy": "registered_style_first_page_render",
        "page_index": 0,
        "reason": "등록공보 감지 → 첫 페이지 렌더 fallback"
    }

def find_representative_drawing_no(page_texts):
    joined = _normalize_text_for_search("\n".join(page_texts))
    patterns = [
        r"【\s*대표도\s*】\s*도\s*([0-9]+)",
        r"대표도\s*[:：\-–—]?\s*도\s*([0-9]+)",
        r"대\s*표\s*도\s*[:：\-–—]?\s*도\s*([0-9]+)",
        r"대표\s*도면\s*[:：\-–—]?\s*도\s*([0-9]+)",
        r"대\s*표\s*도\s*-\s*도\s*([0-9]+)",
    ]
    for pat in patterns:
        m = re.search(pat, joined, flags=re.IGNORECASE)
        if m:
            return int(m.group(1))
    return None

def page_has_meaningful_graphics(page):
    try:
        if page.get_images(full=True):
            return True
    except Exception:
        pass

    try:
        drawings = page.get_drawings()
        if drawings:
            return True
    except Exception:
        pass

    try:
        mat = fitz.Matrix(2, 2)
        pix = page.get_pixmap(matrix=mat, alpha=False)
        pil = Image.open(io.BytesIO(pix.tobytes("png"))).convert("L")
        bbox = ImageChops.invert(pil).getbbox()
        if bbox:
            bw, bh = bbox[2] - bbox[0], bbox[3] - bbox[1]
            if bw * bh > 20000:
                return True
    except Exception:
        pass

    return False

def _drawing_heading_patterns(drawing_no: int):
    n = int(drawing_no)
    return [
        f"【도 {n}】", f"【도{n}】", f"[도 {n}]", f"[도{n}]",
        f"도 {n}", f"도{n}",
    ]

def find_drawing_candidate_pages(doc, page_texts, drawing_no: int):
    pats = [
        rf"【\s*도\s*{drawing_no}\s*】",
        rf"\[\s*도\s*{drawing_no}\s*\]",
        rf"도\s*{drawing_no}\s*은",
        rf"도\s*{drawing_no}\s*를",
        rf"도\s*{drawing_no}\s*에",
        rf"도\s*{drawing_no}\b",
    ]

    candidates = []
    for i, txt in enumerate(page_texts):
        t = txt or ""
        hit = any(re.search(p, t) for p in pats)
        if not hit:
            continue

        score = 0
        if re.search(rf"【\s*도\s*{drawing_no}\s*】", t):
            score += 25
        if re.search(rf"\[\s*도\s*{drawing_no}\s*\]", t):
            score += 18
        if re.search(rf"도\s*{drawing_no}\s*은", t):
            score += 10
        if page_has_meaningful_graphics(doc[i]):
            score += 30
        score += i * 0.1
        candidates.append((score, i))

    candidates.sort(reverse=True)
    return [idx for _, idx in candidates]

def _find_heading_rect_on_page(page, drawing_no: int):
    queries = _drawing_heading_patterns(drawing_no)
    rects = []
    for q in queries:
        try:
            hits = page.search_for(q)
            if hits:
                rects.extend(hits)
        except Exception:
            continue
    if not rects:
        return None
    rects = sorted(rects, key=lambda r: (r.y0, r.x0))
    return rects[0]

def _find_next_heading_y(page, current_y: float):
    txt = page.get_text("text") or ""
    tokens = re.findall(r"【\s*도\s*\d+\s*】", txt)

    next_heading_rects = []
    for token in tokens:
        try:
            rects = page.search_for(token)
            for r in rects:
                if r.y0 > current_y + 5:
                    next_heading_rects.append(r)
        except Exception:
            pass

    if not next_heading_rects:
        return None

    next_heading_rects = sorted(next_heading_rects, key=lambda r: r.y0)
    return next_heading_rects[0].y0

def extract_drawing_region_from_page(doc, page_index: int, drawing_no: int, dpi=REP_IMAGE_DPI):
    page = doc[page_index]
    heading_rect = _find_heading_rect_on_page(page, drawing_no)

    if heading_rect is not None:
        page_rect = page.rect
        next_y = _find_next_heading_y(page, heading_rect.y1)

        top_y = max(heading_rect.y1 + 6, page_rect.y0)
        bottom_y = page_rect.y1 - 8
        if next_y is not None:
            bottom_y = min(bottom_y, next_y - 8)

        if bottom_y - top_y >= 80:
            clip_rect = fitz.Rect(page_rect.x0 + 5, top_y, page_rect.x1 - 5, bottom_y)

            try:
                pil_clip = render_page_to_pil(
                    page, dpi=dpi, clip_rect=clip_rect,
                    apply_trim=False, apply_footer_cut=True, main_figure_crop=True
                )
                ok, reason_ok = is_reasonable_figure_image(pil_clip)
                if ok:
                    return pil_to_png_stream(pil_clip), {
                        "page_index": page_index,
                        "reason": f"도 {drawing_no} heading 하단 영역 클립 + 본체/캡션 crop",
                    }
            except Exception:
                pass

    try:
        pil_full = render_page_to_pil(
            page, dpi=dpi, clip_rect=None,
            apply_trim=True, apply_footer_cut=True, main_figure_crop=True
        )
        ok, reason_ok = is_reasonable_figure_image(pil_full)
        if ok:
            return pil_to_png_stream(pil_full), {
                "page_index": page_index,
                "reason": f"도 {drawing_no} 전체 페이지 fallback 성공({reason_ok})",
            }
    except Exception:
        pass

    pil_full2 = render_page_to_pil(
        page, dpi=dpi, clip_rect=None,
        apply_trim=True, apply_footer_cut=False, main_figure_crop=False
    )
    return pil_to_png_stream(pil_full2), {
        "page_index": page_index,
        "reason": f"도 {drawing_no} 안전 fallback(약한 crop 적용)",
    }

def find_best_fallback_drawing_page(doc, page_texts):
    candidates = []
    for i, txt in enumerate(page_texts):
        t = txt or ""
        score = 0
        if "【도면】" in t:
            score += 5
        if re.search(r"【\s*도\s*\d+\s*】", t):
            score += 22
        if re.search(r"도\s*\d+\s*은", t):
            score += 10
        if page_has_meaningful_graphics(doc[i]):
            score += 25
        score += i * 0.1

        if score > 0:
            candidates.append((score, i))

    if not candidates:
        return 0
    candidates.sort(reverse=True)
    return candidates[0][1]

def extract_representative_from_pdf_unified(pdf_path, dpi=REP_IMAGE_DPI):
    doc = fitz.open(pdf_path)
    page_texts = extract_text_all_pages(doc)

    rep_no = find_representative_drawing_no(page_texts)

    if rep_no is not None:
        candidate_pages = find_drawing_candidate_pages(doc, page_texts, rep_no)
        for page_idx in candidate_pages:
            try:
                stream, info = extract_drawing_region_from_page(doc, page_idx, rep_no, dpi=dpi)
                info["representative_drawing_no"] = rep_no
                info["strategy"] = "representative_drawing_no_match"
                return stream, info
            except Exception:
                continue

    fallback_idx = find_best_fallback_drawing_page(doc, page_texts)
    page = doc[fallback_idx]

    try:
        pil = render_page_to_pil(
            page, dpi=dpi, clip_rect=None,
            apply_trim=True, apply_footer_cut=True, main_figure_crop=True
        )
        ok, reason_ok = is_reasonable_figure_image(pil)
        if ok:
            return pil_to_png_stream(pil), {
                "representative_drawing_no": rep_no,
                "page_index": fallback_idx,
                "strategy": "drawing_page_fallback",
                "reason": f"대표도 번호 미확인 또는 추출 실패 → 도면/그래픽 기반 fallback({reason_ok})"
            }
    except Exception:
        pass

    pil2 = render_page_to_pil(
        page, dpi=dpi, clip_rect=None,
        apply_trim=True, apply_footer_cut=False, main_figure_crop=False
    )
    return pil_to_png_stream(pil2), {
        "representative_drawing_no": rep_no,
        "page_index": fallback_idx,
        "strategy": "drawing_page_fallback_safe",
        "reason": "강한 crop 결과 부적합 → 안전 fallback"
    }

def extract_representative_image_by_meta(pdf_path, meta=None, dpi=REP_IMAGE_DPI):
    meta = meta or {}
    has_reg = bool((meta.get("reg_no") or "").strip() or (meta.get("reg_date") or "").strip())

    if has_reg:
        return extract_representative_from_pdf_registered_style(pdf_path, dpi=dpi)

    return extract_representative_from_pdf_unified(pdf_path, dpi=dpi)

def insert_representative_image(prs, pdf_path, log, meta=None):
    placeholder = "{{[대표이미지]}}"
    slide, shape = find_placeholder_shape(prs, placeholder)
    if not shape:
        log.append("- {{[대표이미지]}}: 템플릿에서 플레이스홀더 미발견 → 스킵")
        return

    try:
        img_stream, info = extract_representative_image_by_meta(pdf_path, meta=meta, dpi=REP_IMAGE_DPI)
        log.append(
            f"- 대표이미지 선택: strategy={info.get('strategy')} / "
            f"page_index={info.get('page_index')} / "
            f"reason={info.get('reason')}"
        )
    except Exception as e:
        log.append(f"- {{[대표이미지]}}: PDF 추출 실패 → 스킵 ({type(e).__name__}: {_shorten_err(e)})")
        return

    left, top, width, height = shape.left, shape.top, shape.width, shape.height
    remove_shape(slide, shape)
    slide.shapes.add_picture(img_stream, left, top, width=width, height=height)
    log.append("- {{[대표이미지]}}: 삽입 완료")

# =========================================================
# 9-1) PPT 강조 처리
# =========================================================
def apply_highlight_in_paragraph(paragraph, log):
    runs = list(paragraph.runs)
    if not runs:
        return

    full = "".join(r.text for r in runs)
    if "⟦" not in full:
        return

    m = _TAGGED_MARK_RE.search(full)
    if not m:
        for r in runs:
            r.text = r.text.replace(HMARK_L, "").replace(HMARK_R, "")
        return

    tag, inner = m.group(1), m.group(2)
    color = COLOR_OLD if tag == "G" else COLOR_NEW

    marked = f"⟦{tag}:{inner}⟧"
    if marked not in full:
        for r in runs:
            r.text = strip_markers(r.text)
        return

    clean_full = full.replace(marked, inner, 1)

    if len(runs) == 1:
        runs[0].text = clean_full
        runs[0].font.bold = True
        runs[0].font.color.rgb = color
        return

    inner_start = clean_full.find(inner)
    if inner_start == -1:
        runs[0].text = clean_full
        runs[0].font.bold = True
        runs[0].font.color.rgb = color
        for j in range(1, len(runs)):
            runs[j].text = ""
        return

    inner_end = inner_start + len(inner)
    prefix = clean_full[:inner_start]
    target = clean_full[inner_start:inner_end]
    suffix = clean_full[inner_end:]

    for r in runs:
        r.text = ""

    runs[0].text = prefix
    if len(runs) >= 3:
        runs[1].text = target
        runs[1].font.bold = True
        runs[1].font.color.rgb = color
        runs[2].text = suffix
    else:
        runs[1].text = target + suffix
        runs[1].font.bold = True
        runs[1].font.color.rgb = color

def apply_highlight_everywhere(prs, log):
    for slide in prs.slides:
        for shape in _iter_shapes_recursive(slide.shapes):
            if getattr(shape, "has_text_frame", False):
                for p in shape.text_frame.paragraphs:
                    apply_highlight_in_paragraph(p, log)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        for p in cell.text_frame.paragraphs:
                            apply_highlight_in_paragraph(p, log)

# =========================================================
# 10) mapping
# =========================================================
def payload_to_mapping(p: dict):
    meta = p.get("meta", {}) if isinstance(p.get("meta", {}), dict) else {}

    app_no = meta.get("app_no", "")
    reg_no = meta.get("reg_no", "")
    app_date = meta.get("app_date", "")
    reg_date = meta.get("reg_date", "")

    app_reg_no = meta.get("app_reg_no_display", "")
    app_reg_date = meta.get("app_reg_date_display", "")
    invention_title = clean_korean_title_only(meta.get("invention_title",""))

    return {
        "{{발명의 명칭}}": invention_title,
        "{{ 발명의 명칭 }}": invention_title,

        "{{출원(등록)번호}}": app_reg_no,
        "{{ 출원(등록)번호 }}": app_reg_no,
        "{{출원(등록)일자}}": app_reg_date,
        "{{ 출원(등록)일자 }}": app_reg_date,

        "{{출원번호}}": app_no,
        "{{ 출원번호 }}": app_no,
        "{{등록번호}}": reg_no,
        "{{ 등록번호 }}": reg_no,
        "{{출원일자}}": app_date,
        "{{ 출원일자 }}": app_date,
        "{{등록일자}}": reg_date,
        "{{ 등록일자 }}": reg_date,

        "{{제목(쉬운 기술카피라이팅)}}": p["tech_name"],
        "{{기술분류}}": p["tech_class"],
        "{{기술명 요약}}": p["tech_summary"],

        "{{적용분야 1}}": p["applications"][0]["title"],
        "{{적용분야 2}}": p["applications"][1]["title"],
        "{{적용분야 3}}": p["applications"][2]["title"],

        "{{본문 1}}": p["overview"][0],
        "{{본문 2}}": p["overview"][1],
        "{{본문 3}}": p["overview"][2],

        "{{기존기술 1}}": p["diff_old"][0],
        "{{기존기술 2}}": p["diff_old"][1],
        "{{기술차별성 1}}": p["diff_new"][0],
        "{{기술차별성 2}}": p["diff_new"][1],
        "{{기술적한계 1}}": p["limits_old"][0],
        "{{기술적한계 2}}": p["limits_old"][1],
        "{{기술적우위 1}}": p["advantages_new"][0],
        "{{기술적우위 2}}": p["advantages_new"][1],
    }

# =========================================================
# 11) 오류 리포트 txt 생성
# =========================================================
def build_error_report_file(status_msg: str, readable: str, payload: dict, log_text: str, tb_text: str = ""):
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    out_path = os.path.join(tempfile.gettempdir(), f"smk_error_report_{ts}.txt")

    sections = []
    sections.append("========================================")
    sections.append("SMK Factory 오류 리포트")
    sections.append("========================================")
    sections.append("")
    sections.append("[상태]")
    sections.append(status_msg or "")
    sections.append("")

    sections.append("[생성된 기술소개자료 텍스트(가~사)]")
    sections.append(readable or "(없음)")
    sections.append("")

    sections.append("[정제 JSON]")
    try:
        sections.append(json.dumps(payload or {}, ensure_ascii=False, indent=2))
    except Exception:
        sections.append("(JSON 직렬화 실패)")
    sections.append("")

    sections.append("[작업 로그]")
    sections.append(log_text or "(로그 없음)")
    sections.append("")

    if tb_text:
        sections.append("[Traceback]")
        sections.append(tb_text)

    with open(out_path, "w", encoding="utf-8") as f:
        f.write("\n".join(sections))

    return out_path

# =========================================================
# 12) 통합 실행 (Streamlit 용도)
# =========================================================
def run_all_st(doc_file, ppt_template_file, page_limit, temperature, log_callback):
    log = []
    readable = ""
    payload = {}

    def add_log(msg):
        log.append(msg)
        log_callback(msg)

    try:
        t0 = time.time()
        add_log("▶ [1/8] 입력 파일 확인 중...")

        file_ext = os.path.splitext(doc_file.name)[1].lower() if hasattr(doc_file, "name") else ".pdf"
        doc_path = get_uploaded_path(doc_file, file_ext)
        if not doc_path:
            err_file = build_error_report_file("문서 업로드/경로 인식 실패", readable, payload, "\n".join(log))
            add_log("❌ 문서 업로드/경로 인식 실패")
            return err_file, True

        template_path = get_uploaded_path(ppt_template_file, ".pptx")
        if not template_path:
            add_log("▶ 기본 템플릿 다운로드 중...")
            if TEMPLATE_PPTX_URL:
                template_path = download_file_to_temp(TEMPLATE_PPTX_URL, ".pptx")
                add_log("- 템플릿 다운로드 완료")
            else:
                err_file = build_error_report_file("PPT 템플릿 업로드 필요함.", readable, payload, "\n".join(log))
                add_log("❌ PPT 템플릿 업로드 필요함")
                return err_file, True

        add_log(f"[모델 고정] 텍스트={TEXT_MODEL_FIXED} / 이미지={IMAGE_MODEL_FIXED}")
        add_log(f"- temperature: {temperature}")
        add_log(f"▶ [2/8] 문서 텍스트 추출 중... (포맷: {file_ext.upper()})")
        
        if file_ext == ".docx":
            patent_text = extract_docx_text(doc_path)
        else:
            patent_text = extract_pdf_text(doc_path, max_pages=int(page_limit))

        if not patent_text:
            err_file = build_error_report_file("텍스트 추출 실패", readable, payload, "\n".join(log))
            add_log("❌ 텍스트 추출 실패함 (스캔된 이미지 파일이거나 빈 파일일 수 있습니다).")
            return err_file, True

        add_log("▶ [3/8] 문서 메타데이터 분석 중...")
        meta = parse_patent_metadata(patent_text)
        if not isinstance(meta, dict):
            meta = {"invention_title": "", "app_no": "", "reg_no": "", "app_date": "", "reg_date": "", "app_reg_no_display": "", "app_reg_date_display": ""}
        
        meta["invention_title"] = clean_korean_title_only(meta.get("invention_title",""))
        add_log(f"- 출원번호='{meta.get('app_no','')}', 등록번호='{meta.get('reg_no','')}'")

        add_log("▶ [4/8] OpenAI 호출 및 기술소개 JSON 생성 중 (시간이 소요됩니다)...")
        payload = call_openai_json(patent_text, temperature, log=log, meta=meta)
        payload["meta"] = meta
        add_log(f"- 쉬운 기술제목 생성값: '{payload.get('tech_name','')}'")

        add_log("▶ [5/8] 텍스트 구조화 및 PPT 템플릿 로드 중...")
        readable = payload_to_readable_text(payload)
        prs = Presentation(template_path)
        mapping = payload_to_mapping(payload)
        replace_counts = replace_placeholders_everywhere(prs, mapping)

        remain = scan_remaining_placeholders_recursive(prs)
        check_placeholders_status(mapping, replace_counts, remain, log, ["{{발명의 명칭}}"])

        add_log("▶ [6/8] 아이콘 생성 및 삽입 중...")
        icon_streams = generate_three_icons(payload, log)
        replace_icons_with_streams(prs, icon_streams, log)

        add_log("▶ [7/8] 대표이미지 추출 및 삽입 중...")
        if file_ext == ".docx":
            add_log("- Word 문서가 업로드되어 대표이미지 자동 추출 기능을 건너뜁니다.")
        else:
            insert_representative_image(prs, doc_path, log, meta=meta)

        add_log("▶ [8/8] 강조 처리 및 최종 저장 준비 중...")
        apply_highlight_everywhere(prs, log)

        ts = datetime.now().strftime("%Y%m%d_%H%M%S")
        out_path = os.path.join(tempfile.gettempdir(), f"output_{ts}.pptx")
        prs.save(out_path)

        add_log(f"\n✅ 총 소요시간: {time.time() - t0:.2f}s")
        return out_path, False

    except Exception:
        tb = traceback.format_exc()
        log_text = "\n".join(log) + "\n\n" + tb
        err_file = build_error_report_file("내부 예외 발생", readable, payload, log_text, tb_text=tb)
        add_log("\n❌ [오류 발생] 작업을 중단하고 에러 리포트를 출력합니다.")
        return err_file, True


# =========================================================
# 13) 외부 호출용 함수 (app.py 연동용)
# =========================================================
def run_smk(spec_file, ppt_template, target_corp=None, ir_data=None, business_status=None):
    """
    app.py에서 호출하는 SMK 생성 엔트리포인트.
    진행 상태 로그는 숨기고, 완료 시 다운로드 버튼만 노출합니다.
    """
    import streamlit as st

    # 화면에 로그를 찍지 않도록 더미(Dummy) 콜백 함수 사용
    def silent_logger(msg):
        pass

    # 생성 진행 중 표시될 로딩 스피너
    with st.spinner("🤖 SMK(기술소개서)를 분석 및 조립 중입니다. 잠시만 기다려주세요..."):
        # run_all_st 함수 호출 (기본값: 25페이지 추출, temperature 0.2)
        out_path, is_error = run_all_st(
            doc_file=spec_file, 
            ppt_template_file=ppt_template, 
            page_limit=25, 
            temperature=0.2, 
            log_callback=silent_logger
        )

    # 결과 출력 및 다운로드 버튼 생성
    if is_error:
        st.error("❌ SMK 생성 중 오류가 발생했습니다. 아래 리포트를 확인해 주세요.")
        with open(out_path, "rb") as file:
            st.download_button(
                label="📥 오류 리포트 다운로드 (.txt)",
                data=file,
                file_name=out_path.split("/")[-1],
                mime="text/plain"
            )
    else:
        st.success("🎉 성공적으로 SMK가 생성되었습니다!")
        with open(out_path, "rb") as file:
            st.download_button(
                label="📥 SMK 최종본 다운로드 (.pptx)",
                data=file,
                file_name="PNU_SMK_Final.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                type="primary"
            )